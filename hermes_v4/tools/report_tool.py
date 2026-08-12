"""Report generation tool for Hermes V4.

Given a topic, asks the LLM for a structured outline (title + sections
+ bullets) and renders it as a PPTX and/or PDF file. Files are handed
back via ToolResult.metadata["file_paths"]; the Telegram layer sends
them as document attachments.
"""

from __future__ import annotations

import asyncio
import json
import logging
import pathlib
import re
import shutil
import tempfile
import time
from typing import Any

from hermes_v4.config.settings import get_settings
from hermes_v4.core.base import Tool, ToolResult
from hermes_v4.llm.provider import LLMProvider
from hermes_v4.prompts import (
    DOCUMENT_SUMMARY_SYSTEM_PROMPT,
    DOCUMENT_TRANSLATE_SYSTEM_PROMPT,
    PPT_TYPE_GUIDE,
    RESEARCH_EXPERT_SYSTEM_PROMPT,
)

logger = logging.getLogger(__name__)

_VALID_FORMATS = {"pptx", "pdf", "both"}

# Formats build_from_document() can extract text from via markitdown.
CONVERTIBLE_DOC_SUFFIXES = {".docx", ".pdf", ".pptx", ".md", ".txt"}

# Above this length, the document is summarized in chunks (map) before
# being handed to outline generation — a flat truncation would silently
# drop whatever didn't fit in the model's context window.
_DOC_CHUNK_CHARS = 6000

# Ollama runs with a small default context window regardless of what the
# model itself supports (observed: a single _DOC_CHUNK_CHARS chunk got cut
# off mid-generation, done_reason "length", after only ~1300 output tokens)
# — every LLM call that processes a full chunk (summarize/translate/PPTX
# doc-to-brief) needs this raised explicitly via num_ctx. Ignored by
# non-Ollama providers, which size context server-side.
_DOC_CHUNK_NUM_CTX = 32768


def _slugify(text: str) -> str:
    slug = re.sub(r"[^\w\s-]", "", text, flags=re.UNICODE).strip()
    slug = re.sub(r"\s+", "_", slug)
    return slug[:60] or "report"


def _is_valid_pptx(path: pathlib.Path) -> bool:
    """A pptx is a zip archive — cheaply guards against rescuing a file
    that was still mid-write when the harness's timeout killed the
    subprocess (a partial zip beats risking a corrupt file being sent)."""
    import zipfile

    try:
        with zipfile.ZipFile(path) as zf:
            return "[Content_Types].xml" in zf.namelist()
    except (zipfile.BadZipFile, OSError):
        return False


def _strip_json_fence(text: str) -> str:
    """Strips a ```json ... ``` (or bare ```) wrapper Claude sometimes adds
    despite being asked for pure JSON."""
    text = text.strip()
    text = re.sub(r"^```[a-zA-Z]*\n?", "", text)
    text = re.sub(r"\n?```$", "", text)
    return text.strip()


class ReportTool(Tool):
    """Generates a PPTX/PDF report for a given topic."""

    name = "generate_report"
    description = (
        "사용자가 발표자료(PPTX)나 문서(PDF) '파일'을 명시적으로 요청했을 때만 사용하세요 — "
        "예: '보고서 만들어줘', 'PPT로 정리해줘', 'PDF로 만들어줘', '슬라이드로 만들어줘'. "
        "주제만 언급하거나 단순히 정보를 물어보는 질문(예: '베샤멜 소스', '~에 대해 알려줘', "
        "'~ 레시피 보내줘')에는 절대 사용하지 마세요 — 그런 경우는 rag나 직접 답변으로 처리해야 "
        "합니다. 호출되면 최신 웹 검색 결과로 리서치한 뒤 파일을 생성합니다 (리서치를 직접 "
        "수행하므로 별도의 web_search 단계를 먼저 호출할 필요는 없습니다)."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "topic": {"type": "string"},
            "format": {"type": "string", "enum": sorted(_VALID_FORMATS)},
        },
        "required": ["topic"],
    }
    output_schema = {
        "type": "object",
        "properties": {
            "summary": {"type": "string"},
            "file_paths": {"type": "array", "items": {"type": "string"}},
        },
    }

    def __init__(
        self,
        llm: LLMProvider,
        web_search_tool: Tool | None = None,
        claude_code_tool: Tool | None = None,
        output_dir: str | pathlib.Path | None = None,
        model: str | None = None,
    ) -> None:
        settings = get_settings()
        self.llm = llm
        self.web_search_tool = web_search_tool
        self.claude_code_tool = claude_code_tool
        self.output_dir = pathlib.Path(output_dir or settings.REPORTS_DIR)
        self.num_sections = settings.REPORT_SECTIONS
        self.bullets_per_section = settings.REPORT_BULLETS_PER_SECTION
        self.num_research_queries = settings.REPORT_RESEARCH_QUERIES
        # Report generation is a one-off, latency-tolerant request (unlike
        # interactive chat), so it's worth spending a bigger/slower local
        # model here for noticeably richer output. Empty = use the LLM
        # provider's own default model.
        self.model = model if model is not None else settings.REPORT_LLM_MODEL or None

    async def validate_input(self, input: dict[str, Any]) -> list[str]:
        errors = []
        topic = input.get("topic")
        if not topic or not str(topic).strip():
            errors.append("'topic' is required and must be a non-empty string")
        fmt = input.get("format", "both")
        if fmt not in _VALID_FORMATS:
            errors.append(f"'format' must be one of {sorted(_VALID_FORMATS)}")
        return errors

    async def execute(self, input: dict[str, Any]) -> ToolResult:
        errors = await self.validate_input(input)
        if errors:
            return ToolResult.fail("; ".join(errors))

        topic = str(input["topic"]).strip()
        fmt = input.get("format", "both")

        start = time.monotonic()
        research, research_fell_back = await self._research(topic)
        try:
            outline, outline_fell_back = await self._generate_outline(topic, research)
        except Exception as exc:
            return ToolResult.fail(f"Failed to generate report outline: {exc}")

        return await self._render(
            outline,
            fallback_title=topic,
            fmt=fmt,
            start=start,
            content_fallback=research_fell_back or outline_fell_back,
        )

    async def build_from_document(
        self, source_path: str | pathlib.Path, fmt: str = "pptx"
    ) -> ToolResult:
        """Converts an uploaded document (docx/pdf/pptx/md/txt) into a
        report, grounded in the document's own content instead of web
        research. Used by the Telegram document handler, not the Planner —
        this is triggered directly by a file upload + conversion request,
        not a text-only tool call."""
        if fmt not in _VALID_FORMATS:
            return ToolResult.fail(f"'format' must be one of {sorted(_VALID_FORMATS)}")

        source_path = pathlib.Path(source_path)
        start = time.monotonic()
        try:
            doc_title, brief = await self._document_to_brief(source_path)
        except Exception as exc:
            return ToolResult.fail(f"문서에서 내용을 추출하지 못했습니다: {exc}")

        try:
            outline, outline_fell_back = await self._generate_outline(
                doc_title, brief, context_label="원본 문서 내용"
            )
        except Exception as exc:
            return ToolResult.fail(f"Failed to generate report outline: {exc}")

        return await self._render(
            outline, fallback_title=doc_title, fmt=fmt, start=start, content_fallback=outline_fell_back
        )

    async def summarize_document(self, source_path: str | pathlib.Path) -> ToolResult:
        """Summarizes an uploaded document (docx/pdf/pptx/md/txt) as chat
        text rather than rendering a file — used by the Telegram document
        handler when the upload caption asks for a summary directly."""
        source_path = pathlib.Path(source_path)
        start = time.monotonic()
        try:
            text = (await self._extract_document_text(source_path)).strip()
        except Exception as exc:
            return ToolResult.fail(f"문서에서 내용을 추출하지 못했습니다: {exc}")
        if not text:
            return ToolResult.fail("문서에서 추출할 수 있는 텍스트가 없습니다.")

        doc_title = source_path.stem
        try:
            if len(text) <= _DOC_CHUNK_CHARS:
                summary = await self._summarize_text(doc_title, text)
            else:
                chunks = [text[i : i + _DOC_CHUNK_CHARS] for i in range(0, len(text), _DOC_CHUNK_CHARS)]
                chunk_summaries = await asyncio.gather(
                    *(self._condense_chunk_for_summary(doc_title, chunk) for chunk in chunks)
                )
                # A second pass over the concatenated per-chunk summaries so
                # the result reads as one coherent summary instead of N
                # disjoint fragments stapled together.
                summary = await self._summarize_text(doc_title, "\n\n".join(chunk_summaries))
        except Exception as exc:
            return ToolResult.fail(f"요약 생성에 실패했습니다: {exc}")

        duration_ms = int((time.monotonic() - start) * 1000)
        return ToolResult(
            success=True,
            output=f"'{doc_title}' 요약:\n\n{summary}",
            metadata={"title": doc_title},
            duration_ms=duration_ms,
        )

    async def translate_document(self, source_path: str | pathlib.Path, target_language: str = "영어") -> ToolResult:
        """Translates an uploaded document (docx/pdf/pptx/md/txt) as chat
        text. Long documents are split into _DOC_CHUNK_CHARS chunks and
        translated in parallel — cheaper and faster than doing it
        sequentially, at the cost of possible terminology drift between
        chunks (no shared context carried across them)."""
        source_path = pathlib.Path(source_path)
        start = time.monotonic()
        try:
            text = (await self._extract_document_text(source_path)).strip()
        except Exception as exc:
            return ToolResult.fail(f"문서에서 내용을 추출하지 못했습니다: {exc}")
        if not text:
            return ToolResult.fail("문서에서 추출할 수 있는 텍스트가 없습니다.")

        doc_title = source_path.stem
        chunks = [text[i : i + _DOC_CHUNK_CHARS] for i in range(0, len(text), _DOC_CHUNK_CHARS)] or [text]
        try:
            translated_chunks = await asyncio.gather(
                *(self._translate_chunk(chunk, target_language) for chunk in chunks)
            )
        except Exception as exc:
            return ToolResult.fail(f"번역에 실패했습니다: {exc}")

        duration_ms = int((time.monotonic() - start) * 1000)
        return ToolResult(
            success=True,
            output=f"'{doc_title}' 번역 ({target_language}):\n\n" + "\n\n".join(translated_chunks),
            metadata={"title": doc_title, "target_language": target_language},
            duration_ms=duration_ms,
        )

    async def _condense_chunk_for_summary(self, doc_title: str, chunk: str) -> str:
        """Condenses one chunk of a long document for summarize_document().

        Deliberately separate from _summarize_document_chunk (used by the
        PPTX/report pipeline, whose prompt explicitly asks for
        presentation-ready material) — reusing that one leaked "발표자료
        개요 작성에 쓸 수 있도록" framing into plain-summary output (PPT-style
        headers/tables even when the user just asked for a summary).
        """
        prompt = (
            f'다음은 문서 "{doc_title}"의 일부 내용입니다. 핵심 주장, 근거, 수치, 결론만 '
            "간결하게 정리하세요 (원문을 그대로 옮기지 말고 요약):\n\n"
            f"{chunk}"
        )
        completion = await self.llm.generate(
            [
                {"role": "system", "content": DOCUMENT_SUMMARY_SYSTEM_PROMPT},
                {"role": "user", "content": prompt},
            ],
            model=self.model,
            num_ctx=_DOC_CHUNK_NUM_CTX,
        )
        return completion.content

    async def _summarize_text(self, doc_title: str, text: str) -> str:
        prompt = (
            f'다음은 문서 "{doc_title}"의 전체 내용입니다. 핵심을 놓치지 않으면서 '
            f"간결하고 읽기 쉬운 요약문을 한국어로 작성하세요:\n\n{text}"
        )
        completion = await self.llm.generate(
            [
                {"role": "system", "content": DOCUMENT_SUMMARY_SYSTEM_PROMPT},
                {"role": "user", "content": prompt},
            ],
            model=self.model,
            num_ctx=_DOC_CHUNK_NUM_CTX,
        )
        return completion.content

    async def _translate_chunk(self, chunk: str, target_language: str) -> str:
        prompt = f"다음 텍스트를 {target_language}로 번역하세요:\n\n{chunk}"
        completion = await self.llm.generate(
            [
                {"role": "system", "content": DOCUMENT_TRANSLATE_SYSTEM_PROMPT},
                {"role": "user", "content": prompt},
            ],
            model=self.model,
            # Translation (unlike summarization) expands rather than
            # compresses the input, so it also needs more headroom above
            # LLM_MAX_TOKENS (4096) than the num_ctx fix alone provides.
            max_tokens=16000,
            num_ctx=_DOC_CHUNK_NUM_CTX,
        )
        return completion.content

    async def _render(
        self,
        outline: dict[str, Any],
        fallback_title: str,
        fmt: str,
        start: float,
        content_fallback: bool = False,
    ) -> ToolResult:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        slug = _slugify(outline.get("title") or fallback_title)
        file_paths = []
        design_fallback = False
        try:
            if fmt in ("pptx", "both"):
                if self.claude_code_tool is not None:
                    path, design_fallback = await self._build_pptx_via_claude_code(outline, slug)
                    file_paths.append(path)
                else:
                    file_paths.append(self._build_pptx(outline, slug))
            if fmt in ("pdf", "both"):
                file_paths.append(self._build_pdf(outline, slug))
        except Exception as exc:
            return ToolResult.fail(f"Failed to render report file: {exc}")

        duration_ms = int((time.monotonic() - start) * 1000)
        summary = (
            f"'{outline.get('title', fallback_title)}' 보고서를 생성했습니다 "
            f"({len(outline.get('sections', []))}개 섹션, {len(file_paths)}개 파일)."
        )
        if content_fallback:
            # Claude Code(구독 기반)를 못 써서 리서치/개요 생성이 로컬
            # 모델(REPORT_LLM_MODEL)로 조용히 강등됐다는 걸 design_fallback과
            # 마찬가지로 알림 — 안 그러면 품질이 낮아진 걸 사용자가 눈치채지
            # 못한 채로 넘어갈 수 있음.
            summary += "\n⚠️ Claude Code를 사용하지 못해서 리서치·개요를 로컬 모델로 생성했어요 — 평소보다 품질이 낮을 수 있어요."
        if design_fallback:
            # AI 디자인 단계가 실패해서 python-pptx 기본 템플릿으로 대체됐다는
            # 걸 조용히 넘기지 않고 알림 — 안 그러면 사용자가 품질이 낮아진 걸
            # 눈치채지 못한 채로 넘어갈 수 있음.
            summary += "\n⚠️ 이번엔 AI 디자인 생성이 실패해서 기본 템플릿으로 만들어졌어요."
        return ToolResult(
            success=True,
            output=summary,
            metadata={
                "file_paths": [str(p) for p in file_paths],
                "title": outline.get("title", fallback_title),
                "design_fallback": design_fallback,
                "content_fallback": content_fallback,
            },
            duration_ms=duration_ms,
        )

    async def _llm_text(self, system: str, prompt: str, json_mode: bool = False) -> tuple[str, bool]:
        """Runs a research/synthesis prompt, preferring Claude Code.

        Claude Code (claude_code_tool) runs with ANTHROPIC_API_KEY unset
        (CLAUDE_CODE_UNSET_API_KEY), so it authenticates via the user's flat-rate
        Claude subscription instead of the metered API — free at the margin,
        unlike adding a real Anthropic API provider here would be. Falls back to
        the local LLMProvider if Claude Code is unavailable or fails, so a CLI
        hiccup never fails the whole report.

        Returns (text, fell_back) — fell_back is True when the local model had
        to be used instead of Claude Code, so callers can surface that to the
        user instead of silently shipping lower-quality content.
        """
        if self.claude_code_tool is not None:
            task = system
            if json_mode:
                task += "\n\n반드시 순수 JSON 텍스트만 응답하세요. 마크다운 코드블록이나 설명을 덧붙이지 마세요."
            task += f"\n\n{prompt}"
            try:
                result = await self.claude_code_tool.execute({"task": task, "purpose": "report_research"})
            except Exception:
                logger.exception("Claude Code text generation failed; falling back to local model")
                result = None
            if result is not None and result.success and result.output:
                text = str(result.output)
                return (_strip_json_fence(text) if json_mode else text), False
            logger.warning("Claude Code text generation unavailable; falling back to local model")

        completion = await self.llm.generate(
            [
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
            model=self.model,
            **({"response_format": "json"} if json_mode else {}),
        )
        return completion.content, True

    async def _generate_search_queries(self, topic: str) -> tuple[list[str], bool]:
        """Fans a topic out into several distinct search angles (background,
        latest developments, stats, counterpoints, examples) instead of one
        flat query — mirrors how a human researcher would approach it."""
        prompt = (
            f'다음 주제를 조사하기 위한 서로 다른 검색어를 {self.num_research_queries}개 '
            f'만드세요: "{topic}"\n\n'
            "각 검색어는 서로 다른 각도(예: 배경/정의, 최신 동향·뉴스, 통계·수치 데이터, "
            "쟁점·반론, 실제 사례)를 다루도록 하고, 실제 검색엔진에 넣을 법한 자연스럽고 "
            "구체적인 문구로 작성하세요. 아래 JSON 형식으로만 응답하세요:\n"
            '{"queries": ["검색어1", "검색어2", ...]}'
        )
        fell_back = False
        try:
            content, fell_back = await self._llm_text(RESEARCH_EXPERT_SYSTEM_PROMPT, prompt, json_mode=True)
            data = json.loads(content)
            queries = [str(q).strip() for q in data.get("queries", []) if str(q).strip()]
        except Exception:
            logger.exception("Failed to generate research sub-queries; falling back to topic itself")
            queries = []
        return (queries[: self.num_research_queries] if queries else [topic]), fell_back

    async def _research(self, topic: str) -> tuple[str, bool]:
        """Best-effort web research to ground the outline. Never fails the
        whole report — if search is unavailable or errors, falls back to
        the LLM's own knowledge (logged, not raised).

        Runs several angled searches in parallel (instead of one flat
        query) and has the local model synthesize them into a research
        brief, so the outline is grounded in cross-checked, deduplicated
        findings rather than raw search snippets.

        Returns (research_brief, fell_back) — fell_back is True if any of
        the LLM calls involved (query generation or synthesis) had to use
        the local model instead of Claude Code.
        """
        if self.web_search_tool is None:
            return "", False

        queries, fell_back = await self._generate_search_queries(topic)

        async def _search_one(query: str) -> str:
            try:
                result = await self.web_search_tool.execute({"query": query})
            except Exception:
                return ""
            if not result.success:
                return ""
            return f"### 검색어: {query}\n{result.output}"

        raw_results = await asyncio.gather(*(_search_one(q) for q in queries))
        raw_results = [r for r in raw_results if r]
        if not raw_results:
            return "", fell_back

        combined = "\n\n".join(raw_results)
        synth_prompt = (
            f'다음은 "{topic}" 주제로 여러 각도에서 검색한 원본 결과입니다. 리서치 '
            "전문가로서 이를 정리하세요: 사실관계를 교차검증하고, 중복된 내용은 합치고, "
            "주제와 무관한 결과는 버리고, 최신 정보와 구체적인 수치·사례를 일반론보다 "
            "우선하세요. 출처 URL이 있다면 관련 문장 끝에 괄호로 남기세요. JSON이 아닌 "
            "서술형 요약문으로 작성하세요.\n\n"
            f"{combined}"
        )
        try:
            text, synth_fell_back = await self._llm_text(RESEARCH_EXPERT_SYSTEM_PROMPT, synth_prompt)
            return text, fell_back or synth_fell_back
        except Exception:
            logger.exception("Research synthesis failed; falling back to raw search results")
            return combined, fell_back

    async def _generate_outline(
        self, topic: str, research: str = "", context_label: str = "웹 검색 결과"
    ) -> tuple[dict[str, Any], bool]:
        research_block = (
            f"\n\n다음은 이 주제와 관련된 {context_label}입니다. 실제로 관련된 내용만 "
            "사실 관계·최신 동향 파악에 참고하고, 무관하거나 엉뚱한 내용은 완전히 무시하세요. "
            "관련 있는 내용이 없다면 이를 참고하지 말고 당신의 기존 지식으로만 작성하세요:\n"
            f"{research}"
            if research
            else ""
        )
        prompt = f"""다음 주제로 발표자료(PPT) 개요를 작성하세요: "{topic}"{research_block}

먼저 아래 10가지 PPT 유형 중 이 주제·목적에 가장 잘 맞는 것을 하나 고르세요. 뚜렷하게 맞는
유형이 없으면 가장 가까운 것을 고르거나 "일반"으로 처리하되, 그 경우 {self.num_sections}개
섹션·섹션당 {self.bullets_per_section}개 bullet을 기본값으로 쓰세요:

{PPT_TYPE_GUIDE}
고른 유형의 전형적인 슬라이드 수·구성 흐름·톤·컬러 방향을 실제로 반영해서 개요를
작성하세요 (섹션 수와 슬라이드당 bullet 개수는 유형에 맞게 자유롭게 정하되, 전체 6~15개
섹션 사이로 하세요). 반드시 한국어로 작성하고, 아래 JSON 형식으로만 응답하세요:

{{"type": "고른 유형", "tone": "톤 설명", "color_hint": "컬러 방향", "title": "보고서 제목", "sections": [{{"heading": "섹션 제목", "bullets": ["내용1", "내용2"]}}]}}"""

        content, fell_back = await self._llm_text(
            "You are a report/presentation outline generator. Respond with JSON only.",
            prompt,
            json_mode=True,
        )
        data = json.loads(content)
        if not isinstance(data, dict) or "sections" not in data:
            raise ValueError(f"LLM returned malformed outline: {content[:200]}")
        return data, fell_back

    async def _extract_document_text(self, path: pathlib.Path) -> str:
        def _convert() -> str:
            from markitdown import MarkItDown

            return MarkItDown().convert(str(path)).text_content or ""

        return await asyncio.to_thread(_convert)

    async def _summarize_document_chunk(self, doc_title: str, chunk: str) -> str:
        prompt = (
            f'다음은 문서 "{doc_title}"의 일부 내용입니다. 발표자료 개요 작성에 쓸 수 있도록 '
            "핵심 주장, 수치, 결론만 간결하게 정리하세요 (원문을 그대로 옮기지 말고 요약):\n\n"
            f"{chunk}"
        )
        completion = await self.llm.generate(
            [
                {"role": "system", "content": RESEARCH_EXPERT_SYSTEM_PROMPT},
                {"role": "user", "content": prompt},
            ],
            model=self.model,
            num_ctx=_DOC_CHUNK_NUM_CTX,
        )
        return completion.content

    async def _document_to_brief(self, path: pathlib.Path) -> tuple[str, str]:
        text = (await self._extract_document_text(path)).strip()
        if not text:
            raise ValueError("문서에서 추출할 수 있는 텍스트가 없습니다.")

        doc_title = path.stem
        if len(text) <= _DOC_CHUNK_CHARS:
            return doc_title, text

        chunks = [text[i : i + _DOC_CHUNK_CHARS] for i in range(0, len(text), _DOC_CHUNK_CHARS)]
        summaries = await asyncio.gather(
            *(self._summarize_document_chunk(doc_title, chunk) for chunk in chunks)
        )
        return doc_title, "\n\n".join(summaries)

    async def _build_pptx_via_claude_code(self, outline: dict[str, Any], slug: str) -> tuple[pathlib.Path, bool]:
        """Delegates slide design to Claude Code, which has Anthropic's own
        "pptx" skill installed (anthropics/skills — see
        ~/.claude/plugins/marketplaces/anthropic-agent-skills). That skill
        writes a pptxgenjs script with real design judgement (palette,
        varied layouts, typography, a visual QA pass rendering every slide
        to an image and checking it) instead of the fixed bullet-list
        template in _build_pptx(). Falls back to that template if Claude
        Code fails for any reason, so a network hiccup or auth issue never
        fails the whole report.
        """
        dest = self.output_dir / f"{slug}.pptx"
        # Claude Code works in an isolated per-call temp dir instead of
        # self.output_dir directly — that directory accumulates every past
        # run's generator scripts (build_*.js) and QA thumbnails
        # (slide-N.jpg) forever, and a fresh `claude -p` invocation reads
        # whatever's already in its cwd, so leftover files from a prior
        # report have caused it to mistake old artifacts for this run's
        # output. Only the final .pptx gets moved into self.output_dir.
        work_dir = pathlib.Path(tempfile.mkdtemp(prefix="pptx_build_", dir=self.output_dir))
        build_dest = work_dir / "output.pptx"
        outline_json = json.dumps(outline, ensure_ascii=False, indent=2)
        task = f"""당신은 기업 임원·투자자 앞에서도 손색없는 슬라이드를 만들어온 시니어
프레젠테이션 디자이너입니다. 콘텐츠의 목적과 청중에 맞는 스타일 판단을 스스로 내리고,
타협 없는 디테일로 마무리합니다.

아래 개요를 바탕으로 전문적이고 시각적으로 매력적인 한국어 PowerPoint(.pptx)
발표자료를 만들어 '{build_dest}' 경로에 저장해주세요.

- 먼저 elite-powerpoint-designer 스킬의 brand-styles.json에서 콘텐츠 톤에 맞는 브랜드
  스타일을 선택하세요 — Tech Keynote, Corporate Professional, Creative Bold, Financial
  Elite, Startup Pitch 외에도 Airbnb Warm, Apple Cinematic, Claude Editorial, IBM
  Carbon, Meta Storefront, Notion Minimal, Nike Bold, Starbucks Café가 있으니 다섯 개로
  좁혀 생각하지 말고 13개 전체 중 콘텐츠 성격에 가장 잘 맞는 것을 고르세요. 그 색상/
  타이포그래피/레이아웃 방향을 pptx 스킬로 실제 파일을 만들 때 반영하세요.
  이 둘은 겹치는 스킬이 아니라 elite-powerpoint-designer가 스타일을, pptx가 실제 생성을
  담당하는 관계이니 둘 다 사용하세요. 개요의 "type"/"tone"/"color_hint" 필드가 있다면
  브랜드 스타일과 색상을 고를 때 반드시 참고하세요 (예: "데이터 분석 보고" 타입이면 숫자를
  큼직하게 강조하는 레이아웃을, "type"에 맞는 톤/컬러를 실제로 반영).
- 섹션 슬라이드마다 서로 다른 레이아웃을 쓰세요 (예: 2단 구성, 아이콘+텍스트 행, 큰 통계
  숫자 강조, 비교 컬럼, 타임라인/프로세스 플로우 등에서 매번 다르게 선택). 모든 섹션에
  똑같은 카드 4개 나열 레이아웃을 반복하지 마세요 — 슬라이드를 나란히 렌더링해서 비교했을 때
  구조가 겹치는 슬라이드가 있으면 안 됩니다.
- "원형 아이콘 배지 + 카드/리스트" 조합(아이콘을 동그라미 안에 넣고 그 옆/아래에 제목+설명을
  나열하는 패턴)은 덱 전체에서 최대 1~2장까지만 쓰세요. 안전한 선택이라고 매 슬라이드에
  변형만 줘서 반복하면(리스트형 vs 그리드형 vs 다크배경형이어도 같은 패턴으로 침) 전체
  덱이 밋밋해 보입니다. 나머지 슬라이드는 실제 차트(막대/라인/도넛), 큰 풀 쿼트, 비대칭
  2단/3단 분할, 사진·일러스트 자리, 프로세스 플로우 등 서로 다른 시각 언어를 적극 쓰세요.
- 무난하고 안전한 스타일로 수렴하지 마세요. elite-powerpoint-designer로 스타일을 고를 때
  콘텐츠가 데이터 중심이 아니라면 Corporate Professional보다 과감한 선택(Creative Bold,
  강한 액센트 컬러, 레이어드 배경 그래픽, 비대칭 구성)을 우선 고려하고, 고른 이유를 스스로
  점검하세요 — "무난해서" 골랐다면 다시 판단하세요.
- 한글 폰트는 "나눔고딕(NanumGothic)"을 사용하세요 — 맑은 고딕은 Windows 전용이라 이
  macOS 렌더링/QA 환경과 사용자 기기 상당수에 없어서, 없는 폰트로 지정하면 뷰어가 알아볼 수
  없는 엉뚱한 장식체로 대체해버려 본문 전체가 읽기 힘든 필기체처럼 나옵니다. 나눔고딕은 이
  환경에 실제로 설치되어 있어 QA 렌더링에서도 지정한 대로 정확히 보입니다.
- **`fontFace: "나눔고딕"`을 제목뿐 아니라 addText()·addTable()·addChart() 등 텍스트가
  들어가는 모든 호출에 예외 없이 개별적으로 명시하세요.** 실제로 자주 발생하는 실패
  패턴은 큰 제목/굵은 글씨에는 나눔고딕을 지정해놓고, 캡션·부제·본문·라벨 같은 작은/
  보통 굵기 텍스트는 fontFace를 빠뜨리거나 브랜드 스타일의 영문 폰트명(SF Pro, IBM Plex
  Sans 등)을 그대로 넣는 것입니다 — 그 영문 폰트들은 한글 글리프가 없어서, 렌더러가 임의로
  다른 한글 폰트로 대체하며 브러시/캘리그래피처럼 보이는 엉뚱한 장식체가 나올 수 있습니다.
  브랜드 스타일의 폰트는 두께감·자간 등 방향성 참고용으로만 쓰고, 실제 fontFace 값은
  본문/캡션/축라벨/범례를 포함해 모든 텍스트에 항상 "나눔고딕"으로 고정하세요. QA 이미지에서
  큰 글씨만 확인하지 말고 작은 캡션·라벨 텍스트도 확대해서 브러시체로 깨지지 않았는지
  반드시 확인하세요.
- pptx 스킬의 타이포그래피 표(제목 36-44pt, 섹션헤더 20-24pt, 본문 14-16pt, 캡션
  10-12pt)와 여백 규칙(최소 0.5" 마진)을 그대로 따르세요 — 임의로 더 작게 줄이지 마세요.
- `addChart()`를 쓸 때 라이브러리 기본 크기 그대로 두지 말고, 그 차트가 슬라이드에서
  가장 큰 시각 요소가 되도록 폭/높이(w/h)를 명시적으로 지정하세요. 차트를 작게 그려놓고
  나머지를 빈 여백으로 채우는 게 가장 흔한 실패 패턴입니다.
- 완료 후 반드시 QA(슬라이드 이미지 렌더링 확인, validate 스크립트)를 거쳐 '{build_dest}'에
  파일이 정상 생성됐는지 확인하세요. 이미지 QA 시 각 슬라이드의 한글 본문이 실제로 고딕체로
  또렷하게 보이는지 확인하고, pptx 스킬의 QA 체크리스트에 있는 "Uneven gaps (large empty
  area in one place, cramped in another)" 항목을 슬라이드마다 빠짐없이 실제로 점검하세요 —
  이 항목을 대충 넘기지 말고, 차트/도형 옆에 불필요하게 큰 빈 공간이 있으면 요소를 키우거나
  재배치해서 다시 렌더링해 재확인하세요.

개요:
{outline_json}
"""
        call_start = time.monotonic()
        try:
            result = await self.claude_code_tool.execute(
                {
                    "task": task,
                    "cwd": str(work_dir),
                    "permission_mode": "acceptEdits",
                    "purpose": "pptx_design",
                    # acceptEdits only auto-approves Write/Edit — every Bash
                    # call the pptx skill's workflow needs (pptxgenjs via
                    # node, npm install as a require() fallback, LibreOffice
                    # conversion, pdftoppm/markitdown for QA, zip/rm while
                    # unpacking a template) still prompts, and headless mode
                    # has no TTY to approve it — so each has to be
                    # pre-approved explicitly or that step silently hangs.
                    "allowed_tools": [
                        "Bash(python3 *)",
                        "Bash(python *)",
                        "Bash(node *)",
                        "Bash(npm *)",
                        "Bash(npx *)",
                        "Bash(pdftoppm *)",
                        "Bash(markitdown *)",
                        "Bash(zip *)",
                        "Bash(unzip *)",
                        "Bash(rm *)",
                        "Bash(fc-list *)",
                    ],
                }
            )
        except Exception:
            logger.exception("Claude Code PPTX design raised an exception; falling back to template")
            result = None
        logger.info(
            "Claude Code PPTX call returned after %.1fs, build_dest.exists()=%s, work_dir=%s",
            time.monotonic() - call_start, build_dest.exists(), sorted(p.name for p in work_dir.iterdir()),
        )

        # The deck is often fully written well before the whole task
        # finishes (QA re-renders, cleanup, etc. run after it) — if the
        # harness's own timeout (CLAUDE_CODE_TIMEOUT) fires during that
        # tail end, result.success is False even though a complete, valid
        # .pptx is already sitting in work_dir. Rescue it instead of
        # discarding a finished deck for the bare-bones template.
        if build_dest.exists() and _is_valid_pptx(build_dest):
            if result is None or not result.success:
                logger.warning(
                    "Claude Code call didn't report success (%s) but '%s' exists and is a valid pptx; using it anyway",
                    result.error if result is not None else "no result", build_dest,
                )
            shutil.move(str(build_dest), str(dest))
            shutil.rmtree(work_dir, ignore_errors=True)
            return dest, False

        if result is not None and result.success:
            # The claude subprocess reporting done doesn't guarantee the
            # file it wrote is visible to us yet — give it a window rather
            # than immediately declaring failure and overwriting a file
            # that's about to land.
            for _ in range(40):
                if build_dest.exists() and _is_valid_pptx(build_dest):
                    shutil.move(str(build_dest), str(dest))
                    shutil.rmtree(work_dir, ignore_errors=True)
                    return dest, False
                await asyncio.sleep(1)
            logger.warning(
                "Still not visible after waiting; work_dir now=%s",
                sorted(p.name for p in work_dir.iterdir()),
            )

        shutil.rmtree(work_dir, ignore_errors=True)

        if result is None:
            logger.warning("Claude Code PPTX design unavailable; falling back to template")
        elif not result.success:
            logger.warning("Claude Code PPTX design failed (%s); falling back to template", result.error)
        else:
            logger.warning(
                "Claude Code reported success but '%s' was not created; falling back to template. "
                "output=%r permission_denials=%r",
                dest, result.output, result.metadata.get("permission_denials"),
            )
        return self._build_pptx(outline, slug), True

    def _build_pptx(self, outline: dict[str, Any], slug: str) -> pathlib.Path:
        from pptx import Presentation

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = outline.get("title", "")
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Hermes Agent 자동 생성 보고서"

        bullet_layout = prs.slide_layouts[1]
        for section in outline.get("sections", []):
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = section.get("heading", "")
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = section.get("bullets", [])
            for i, bullet in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = str(bullet)

        dest = self.output_dir / f"{slug}.pptx"
        prs.save(str(dest))
        return dest

    def _build_pdf(self, outline: dict[str, Any], slug: str) -> pathlib.Path:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer

        registered = set(pdfmetrics.getRegisteredFontNames())
        for font_name in ("HYGothic-Medium", "HYSMyeongJo-Medium"):
            if font_name not in registered:
                pdfmetrics.registerFont(UnicodeCIDFont(font_name))

        styles = getSampleStyleSheet()
        styles["Title"].fontName = "HYGothic-Medium"
        styles["Heading2"].fontName = "HYGothic-Medium"
        styles["Normal"].fontName = "HYSMyeongJo-Medium"

        dest = self.output_dir / f"{slug}.pdf"
        doc = SimpleDocTemplate(str(dest), pagesize=A4)
        story = [Paragraph(outline.get("title", ""), styles["Title"]), Spacer(1, 20)]
        for section in outline.get("sections", []):
            story.append(Paragraph(section.get("heading", ""), styles["Heading2"]))
            bullets = section.get("bullets", [])
            story.append(ListFlowable([ListItem(Paragraph(str(b), styles["Normal"])) for b in bullets]))
            story.append(Spacer(1, 12))
        doc.build(story)
        return dest
